from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from dotenv import load_dotenv
import sys
import os

# Add the parent directory to the Python path
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from agents.agent_manager import AgentManager

load_dotenv()

class PPTGenerator:
    def __init__(self):
        self.agent_manager = AgentManager()
        
    def _apply_title_style(self, title_shape):
        """Apply consistent styling to slide titles"""
        title_format = title_shape.text_frame.paragraphs[0].font
        title_format.name = '微软雅黑'
        title_format.size = Pt(40)
        title_format.bold = True
        title_format.color.rgb = RGBColor(0, 76, 153)  # Dark blue
        
    def _apply_content_style(self, content_frame, is_bullet=True):
        """Apply consistent styling to slide content"""
        for paragraph in content_frame.paragraphs:
            if is_bullet:
                paragraph.level = 0
            font = paragraph.font
            font.name = '微软雅黑'
            font.size = Pt(24)
            font.color.rgb = RGBColor(0, 0, 0)
            
    def _add_title_slide(self, prs, title_text):
        """Add a title slide with proper styling"""
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        subtitle.text = "数学课程"
        
        # Style the title
        title_format = title.text_frame.paragraphs[0].font
        title_format.name = '微软雅黑'
        title_format.size = Pt(54)
        title_format.bold = True
        title_format.color.rgb = RGBColor(0, 76, 153)
        
        # Style the subtitle
        subtitle_format = subtitle.text_frame.paragraphs[0].font
        subtitle_format.name = '微软雅黑'
        subtitle_format.size = Pt(32)
        subtitle_format.color.rgb = RGBColor(89, 89, 89)
        
    def _add_content_slide(self, prs, title_text, content_text):
        """Add a content slide with proper styling"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
        
        # Set title
        title = slide.shapes.title
        title.text = title_text
        self._apply_title_style(title)
        
        # Set content
        content = slide.placeholders[1]
        tf = content.text_frame
        
        # Split content into bullet points if it contains newlines
        points = content_text.strip().split('\n')
        
        # Add each point as a paragraph
        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point.strip()
            p.level = 0 if not point.startswith('  ') else 1
            
        self._apply_content_style(tf)
        
    def _add_diagram_slide(self, prs, title_text, content_text):
        """Add a slide with a basic diagram for mathematical concepts"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
        
        # Add title
        title = slide.shapes.title
        title.text = title_text
        self._apply_title_style(title)
        
        # Add text box for content
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(4)
        height = Inches(4)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.text = content_text
        self._apply_content_style(tf.paragraphs[0], is_bullet=False)
        
        # Add basic shape (e.g., for integral symbol or graph)
        shape_left = Inches(5)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.MATHEMATICAL_INTEGRAL,
            shape_left, top, Inches(4), Inches(4)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 76, 153)
        
    def generate(self, outline):
        """
        Generate PowerPoint slides based on the outline
        
        Args:
            outline (str): Teacher's outline for the lesson
            
        Returns:
            dict: Dictionary containing slide contents and metadata
        """
        # Create presentation object
        prs = Presentation()
        
        # Set slide dimensions to 16:9
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        
        # Generate content using Autogen agents
        print('generating content')
        response = self.agent_manager.generate_content(outline)
        print('ppt generated')
        slides_content = response['ppt_content']['slides']
        
        # Add title slide
        self._add_title_slide(prs, "黎曼积分")
        
        # Create content slides
        for i, slide_content in enumerate(slides_content):
            if "积分" in slide_content['title'] or "定义" in slide_content['title']:
                # Use diagram slide for integral concepts
                self._add_diagram_slide(
                    prs,
                    slide_content['title'],
                    slide_content['content']
                )
            else:
                # Use regular content slide
                self._add_content_slide(
                    prs,
                    slide_content['title'],
                    slide_content['content']
                )
            
        # Save presentation
        output_path = "temp_presentation.pptx"
        prs.save(output_path)
        
        return {
            'slides': slides_content,
            'file_path': output_path
        }
